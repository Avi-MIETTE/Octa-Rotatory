import os
from pptx import Presentation
from docx import Document
from docx.shared import Inches

# ppt file ko doc file mein convert karta hai yeah file

# Folder where your PPT files are stored
input_folder = r"D:\xyz ppt"
# Folder where DOCX files will be saved
output_folder = r"D:\xyz Doc"

os.makedirs(output_folder, exist_ok=True)


def ppt_to_docx(pptx_path, output_path):
    prs = Presentation(pptx_path)
    doc = Document()
    image_count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 = picture
                image = shape.image
                image_bytes = image.blob
                img_name = os.path.join(output_folder, f"temp_{image_count}.png")
                with open(img_name, "wb") as f:
                    f.write(image_bytes)
                # Add image to docx (one below the other)
                doc.add_picture(img_name, width=Inches(5.5))
                doc.add_paragraph()  # Adds space between images
                image_count += 1
                os.remove(img_name)

    doc.save(output_path)
    print(f"Saved: {output_path}")


for file in os.listdir(input_folder):
    if file.lower().endswith(".pptx"):
        ppt_path = os.path.join(input_folder, file)
        docx_name = os.path.splitext(file)[0] + ".docx"
        output_path = os.path.join(output_folder, docx_name)
        ppt_to_docx(ppt_path, output_path)

print("All PPT files converted successfully.")
